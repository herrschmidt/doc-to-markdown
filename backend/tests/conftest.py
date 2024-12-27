import os
import pytest
import logging
from pathlib import Path
from fastapi.testclient import TestClient
from PIL import Image, ImageDraw, ImageFont
from reportlab.pdfgen import canvas
from io import BytesIO
from docx import Document
from app.main import app

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

@pytest.fixture
def test_client():
    """Create a test client for the FastAPI app."""
    return TestClient(app)

@pytest.fixture
def fixtures_dir():
    """Get the path to the test fixtures directory."""
    fixtures_dir = Path(__file__).parent / "fixtures"
    fixtures_dir.mkdir(exist_ok=True)
    return fixtures_dir

@pytest.fixture
def sample_pdf(fixtures_dir):
    """Create a sample PDF file for testing."""
    pdf_path = fixtures_dir / "sample.pdf"
    logger.info(f"Creating sample PDF at: {pdf_path}")
    
    # Always create a new PDF file to ensure it's valid
    logger.info("Creating new PDF file with test content")
    c = canvas.Canvas(str(pdf_path))
    c.setFont("Helvetica", 12)
    c.drawString(100, 750, "Test PDF Document")
    c.drawString(100, 730, "This is a test document created for testing purposes.")
    c.drawString(100, 710, "It contains some text that should be converted to markdown.")
    c.save()
    logger.info(f"PDF file created successfully, size: {pdf_path.stat().st_size} bytes")
    
    return pdf_path

@pytest.fixture
def sample_image(fixtures_dir):
    """Create a sample PNG image for testing."""
    img_path = fixtures_dir / "sample.png"
    logger.info(f"Creating sample PNG at: {img_path}")
    
    # Always create a new PNG file to ensure it's valid
    logger.info("Creating new PNG image with test content")
    img = Image.new('RGB', (1200, 800), color='white')
    d = ImageDraw.Draw(img)
    
    # Try to use a system font
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
        small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
    except:
        font = ImageFont.load_default()
        small_font = ImageFont.load_default()
    
    text = [
        ("Sample Document", font, 50),
        ("", small_font, 120),
        ("This is a test image with multiple lines of text.", small_font, 180),
        ("It includes:", small_font, 240),
        ("• A title", small_font, 300),
        ("• Multiple paragraphs", small_font, 360),
        ("• And bullet points", small_font, 420),
        ("", small_font, 480),
        ("Let's see how well the OCR works!", small_font, 540)
    ]
    
    for line, font_to_use, y_pos in text:
        d.text((50, y_pos), line, fill='black', font=font_to_use)
    
    img.save(img_path, format='PNG')
    logger.info(f"PNG file created successfully, size: {img_path.stat().st_size} bytes")
    
    return img_path

@pytest.fixture
def sample_docx(fixtures_dir):
    """Create a sample DOCX file for testing."""
    docx_path = fixtures_dir / "sample.docx"
    logger.info(f"Creating sample DOCX at: {docx_path}")
    
    # Create a valid DOCX file with some content
    logger.info("Creating new DOCX file with test content")
    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test document created for testing purposes.')
    doc.add_paragraph('It contains some text that should be converted to markdown.')
    doc.save(str(docx_path))
    logger.info(f"DOCX file created successfully, size: {docx_path.stat().st_size} bytes")
    
    return docx_path

@pytest.fixture
def sample_pptx(fixtures_dir):
    """Create a sample PPTX file for testing."""
    from pptx import Presentation
    from pptx.util import Inches
    
    pptx_path = fixtures_dir / "sample.pptx"
    logger.info(f"Creating sample PPTX at: {pptx_path}")
    
    # Create a valid PPTX file with some content
    logger.info("Creating new PPTX file with test content")
    prs = Presentation()
    
    # Add a title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created for testing purposes"
    
    # Add a content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Content Slide"
    body.text = "This is a test slide with some content."
    
    prs.save(str(pptx_path))
    logger.info(f"PPTX file created successfully, size: {pptx_path.stat().st_size} bytes")
    
    return pptx_path

@pytest.fixture
def sample_html(fixtures_dir):
    """Create a sample HTML file for testing."""
    html_path = fixtures_dir / "sample.html"
    logger.info(f"Creating sample HTML at: {html_path}")
    
    # Create a valid HTML file with some content
    logger.info("Creating new HTML file with test content")
    html_content = """<!DOCTYPE html>
<html>
<head>
    <title>Test Document</title>
</head>
<body>
    <h1>Test Document</h1>
    <p>This is a test document created for testing purposes.</p>
    <p>It contains some text that should be converted to markdown.</p>
    <ul>
        <li>First bullet point</li>
        <li>Second bullet point</li>
    </ul>
</body>
</html>"""
    html_path.write_text(html_content)
    logger.info(f"HTML file created successfully, size: {html_path.stat().st_size} bytes")
    
    return html_path

@pytest.fixture
def cleanup_fixtures(fixtures_dir):
    """Clean up test files after tests (disabled for debugging)."""
    yield
    # Cleanup disabled to allow inspection of test files
    logger.info(f"Test files available in: {fixtures_dir}")
    for file in fixtures_dir.glob("*"):
        if file.is_file():
            logger.info(f"  {file.name}: {file.stat().st_size} bytes")
