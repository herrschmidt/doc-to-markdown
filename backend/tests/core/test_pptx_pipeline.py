import pytest
from pathlib import Path
from pptx import Presentation
from app.core.converter import DocumentConverter

@pytest.fixture
def sample_pptx_file(tmp_path):
    # Create a sample PowerPoint presentation
    prs = Presentation()
    
    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Sample Presentation"
    subtitle.text = "Created for testing"
    
    # Add a content slide with bullet points
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Test Slide"
    
    tf = body_shape.text_frame
    tf.text = "First Level"
    p = tf.add_paragraph()
    p.text = "Second Level"
    p.level = 1
    
    # Add speaker notes
    notes_slide = bullet_slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes for testing"
    
    # Save the presentation
    file_path = tmp_path / "sample.pptx"
    prs.save(file_path)
    return file_path

async def test_pptx_conversion(sample_pptx_file):
    converter = DocumentConverter()
    
    # Create a mock UploadFile
    class MockUploadFile:
        def __init__(self, path):
            self.filename = path.name
            self._path = path
        
        async def read(self):
            return self._path.read_bytes()
    
    mock_file = MockUploadFile(sample_pptx_file)
    save_path = Path(sample_pptx_file).parent / "temp.pptx"
    
    try:
        result = await converter.convert(mock_file, save_path)
        
        # Verify the conversion result
        assert result["content"] is not None
        assert "Sample Presentation" in result["content"]
        assert "Test Slide" in result["content"]
        assert "First Level" in result["content"]
        assert "Second Level" in result["content"]
        
        # Verify metadata
        assert result["metadata"]["mime_type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert result["metadata"]["original_file"] == "sample.pptx"
        assert result["metadata"]["file_size"] > 0
        
    finally:
        # Clean up
        if save_path.exists():
            save_path.unlink()